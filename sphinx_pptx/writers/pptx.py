"""Custom docutils writer for plain text."""
from __future__ import annotations

import math
import os
import re
import textwrap
from itertools import chain, groupby
from typing import TYPE_CHECKING, Any, Generator, Iterable, cast

from docutils import nodes, writers
from docutils.nodes import Element, Text
from docutils.utils import column_width

from sphinx import addnodes
from sphinx.locale import _, admonitionlabels
from sphinx.util.docutils import SphinxTranslator

from pptx import Presentation

if TYPE_CHECKING:
    from .builders.text import PptxBuilder





class PptxWriter(writers.Writer):
    supported = ('text',)
    settings_spec = ('No options here.', '', ())
    settings_defaults: dict[str, Any] = {}

    output: str = None

    def __init__(self, builder: PptxBuilder) -> None:
        super().__init__()
        self.builder = builder

    def translate(self) -> None:
        visitor = self.builder.create_translator(self.document, self.builder)
        self.document.walkabout(visitor)
        #self.output = cast(PptxTranslator, visitor).body
        self.prs = cast(PptxTranslator, visitor).prs


class PptxTranslator(SphinxTranslator):
    builder: PptxBuilder

    def __init__(self, document: nodes.document, builder: PptxBuilder) -> None:
        super().__init__(document, builder)

        config_template = self.config.presentation_template
        try:
            self.prs = Presentation(config_template)
        except:
            if config_template is not None:
                import pdb;pdb.set_trace()
                print(f"FAILURE: Could not find {config_template}")
            self.prs = Presentation()
        self.current_slide = None
        self.title_active = False
        self.header_active = False
        self.footer_active = False

    def visit_slide_node(self, node) -> None:
        title_slide_layout = self.prs.slide_layouts[1]
        self.current_slide = self.prs.slides.add_slide(title_slide_layout)

    def depart_slide_node(self, node) -> None:
        self.current_slide = None

    def visit_Text(self, node) -> None:
        if self.footer_active or self.header_active: 
            return
        node_text = node.astext()
        if self.title_active or node_text[0] == '#':
            while node_text[0] == '#':
                node_text = node_text[1:]
            node_text = node_text.strip()
            placeholder_index = 0
        else:
            placeholder_index = 1
        if self.current_slide is not None:
            if self.current_slide.placeholders[placeholder_index].text != '':
                additional_text = "\n" + node_text
            else:
                additional_text = node_text
            self.current_slide.placeholders[placeholder_index].text += additional_text

    def depart_Text(self, *args) -> None:
        pass

    def visit_pseudo_heading_node(self, node) -> None:
        self.title_active = True

    def depart_pseudo_heading_node(self, node) -> None:
        self.title_active = False

    def visit_header_node(self, node) -> None:
        self.header_active = True

    def depart_header_node(self, node) -> None:
        self.header_active = False

    def visit_footer_node(self, node) -> None:
        self.footer_active = True

    def depart_footer_node(self, node) -> None:
        self.footer_active = False

    

no_action_tags = [
    'div_wrap_node',
    'paragraph',
    'document',
    'comment',
    'title',
    'compound',
    'section',
    'inline',
    'reference',
    'listitem',
    'list_item',
    'bullet_list',
    'title',
    'background_node',
    'content_center_node',
    'line',
    'line_block',
    'paragraph_style_node',
    'paragraph_node',
    'content_right_node',
    'template_node',
    'image',
    'block_quote',
    'strong',
    'cite_node',
    'text_quote_node',
    'term',
    'definition',
    'definition_list_item',
    'definition_list',
    'enumerated_list',
]

for tag in no_action_tags:
    def placeholder(self, node):
        pass
    setattr(
        PptxTranslator,
        f'visit_{tag}',
        placeholder,
    )
    setattr(
        PptxTranslator,
        f'depart_{tag}',
        placeholder,
    )
    